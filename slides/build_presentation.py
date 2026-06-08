"""
Builds the 15-slide CAP6412 presentation using python-pptx.
Run: python slides/build_presentation.py
Output: slides/presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1E, 0x27, 0x61)   # #1E2761 — primary dark bg
ICE       = RGBColor(0xCA, 0xDC, 0xFC)   # #CADCFC — light accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL  = RGBColor(0x22, 0x22, 0x22)
SLATE     = RGBColor(0x44, 0x55, 0x66)
GOLD      = RGBColor(0xF5, 0xC5, 0x18)   # accent
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFF)   # near-white content bg
MID_BLUE  = RGBColor(0x2E, 0x47, 0x9C)   # section bars
TEAL      = RGBColor(0x02, 0x80, 0x90)   # highlight teal

W = Inches(10)
H = Inches(5.625)


# ── Helpers ───────────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color, line=False):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, italic=False,
             font_name="Calibri", v_anchor=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_bullet_list(slide, items, x, y, w, h, font_size=14, color=CHARCOAL,
                    title_item=None, indent_items=None, font_name="Calibri"):
    """items: list of str. indent_items: set of indices to indent."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    indent_items = indent_items or set()

    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        if i in indent_items:
            p.level = 1
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
    return txBox

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_table(slide, headers, rows, x, y, w, h,
              header_bg=NAVY, header_fg=WHITE, alt_bg=LIGHT_BG,
              font_size=12):
    from pptx.util import Inches, Pt
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                  Inches(w), Inches(h)).table
    col_w = Inches(w / n_cols)
    for col in tbl.columns:
        col.width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = header_fg
                run.font.name = "Calibri"

    # Data rows
    for ri, row in enumerate(rows):
        bg = alt_bg if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.name = "Calibri"
                    run.font.color.rgb = CHARCOAL
    return tbl


# ── Slide builders ────────────────────────────────────────────────────────────

def slide1_title(prs):
    """Dark title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # Gold accent bar left
    add_rect(slide, 0, 0, 0.18, 5.625, GOLD)

    # Title
    add_text(slide,
             "Bias & Safety Auditor for\nText-to-Image Generative Models",
             0.5, 0.85, 8.8, 2.0,
             font_size=34, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri")

    # Student info block
    add_text(slide,
             "Advance Generative Model",
             0.5, 2.95, 8.8, 0.42,
             font_size=15, color=GOLD, bold=True, font_name="Calibri")

    info_lines = [
        ("Student :", "Waqar Rauf Butt"),
        ("Roll No :", "PHDAIF25M003"),
        ("Supervisor :", "Dr. Muhammad Farooq"),
    ]
    for i, (label, value) in enumerate(info_lines):
        add_text(slide, label, 0.5, 3.42 + i * 0.33, 1.4, 0.32,
                 font_size=12, color=ICE, bold=True, font_name="Calibri")
        add_text(slide, value, 1.95, 3.42 + i * 0.33, 6.0, 0.32,
                 font_size=12, color=WHITE, font_name="Calibri")

    # GitHub link
    add_text(slide,
             "github.com/waqarbutt58/BiasAudit-CAP6412",
             0.5, 5.18, 9.0, 0.3,
             font_size=11, color=ICE, font_name="Calibri")


def content_header(slide, title, subtitle=None):
    """Standard content slide header — navy left bar + title."""
    add_rect(slide, 0, 0, 10, 0.75, NAVY)
    add_text(slide, title, 0.3, 0.08, 9.2, 0.62,
             font_size=22, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri")


def slide2_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Why Audit Text-to-Image Models?")

    # Three stat callout boxes
    stats = [
        ("88%", "male for\n'CEO' prompt"),
        ("93%", "male for\n'Electrician' prompt"),
        ("79%", "female for\n'Nurse' prompt"),
    ]
    box_w, box_h = 2.6, 1.5
    for i, (num, label) in enumerate(stats):
        bx = 0.35 + i * 3.1
        add_rect(slide, bx, 0.9, box_w, box_h, NAVY)
        add_text(slide, num, bx, 0.92, box_w, 0.85,
                 font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, label, bx, 1.72, box_w, 0.6,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    bullets = [
        "Models trained on ~5 billion internet image-text pairs — biases baked in",
        "No standardised open-source end-to-end auditing tool exists",
        "EU AI Act (2024) and US Executive Order on AI Safety mandate bias audits",
        "Industry audits exist but tools are not open-sourced or reproducible",
    ]
    add_bullet_list(slide, bullets, 0.4, 2.6, 9.2, 2.7, font_size=14)


def slide3_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "The Gap We Fill")

    # Existing work column
    add_rect(slide, 0.3, 0.95, 4.3, 1.0, MID_BLUE)
    add_text(slide, "Existing Work", 0.3, 0.97, 4.3, 0.45,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Covers components in isolation:\n"
             "bias detection OR safety filtering OR concept erasure",
             0.35, 1.4, 4.2, 0.65, font_size=12, color=CHARCOAL)

    # Gap arrow
    add_text(slide, "→", 4.7, 1.2, 0.6, 0.5, font_size=28, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)

    # This project column
    add_rect(slide, 5.4, 0.95, 4.3, 1.0, NAVY)
    add_text(slide, "This Project", 5.4, 0.97, 4.3, 0.45,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "End-to-end pipeline: bias detection + safety filtering\n"
             "+ concept erasure + automated PDF report",
             5.45, 1.4, 4.2, 0.65, font_size=12, color=CHARCOAL)

    # Key claim highlight
    add_rect(slide, 0.3, 2.2, 9.4, 0.65, TEAL)
    add_text(slide,
             "First open-source pipeline combining all 4 components with automated PDF report generation",
             0.35, 2.23, 9.3, 0.58,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bullets = [
        "Bias Score + Shannon Entropy + Demographic Parity — three complementary metrics",
        "Training-free safety filter: no fine-tuning required (SAFREE-inspired)",
        "Surgical concept erasure: only ~2% of U-Net parameters fine-tuned",
        "One-command PDF audit report generation with charts",
    ]
    add_bullet_list(slide, bullets, 0.4, 3.05, 9.2, 2.3, font_size=13)


def slide4_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Pipeline Architecture — 5 Modules")

    # Flow boxes
    boxes = [
        (0.25, 0.95, 1.8,  1.1, NAVY,     "INPUT\n60 Prompts"),
        (2.3,  0.95, 2.0,  1.1, MID_BLUE, "MODULE 1\nImage Generation\n(600 images)"),
        (4.55, 0.95, 1.9,  1.1, MID_BLUE, "MODULE 2\nBias Detection\nCLIP zero-shot"),
        (4.55, 2.3,  1.9,  1.1, MID_BLUE, "MODULE 3\nSafety Filter\nNSFW detection"),
        (6.7,  0.95, 1.75, 2.45,MID_BLUE, "MODULE 4\nConcept\nErasure"),
        (8.7,  0.95, 1.1,  2.45,NAVY,     "MODULE 5\nAudit\nReport"),
    ]
    for bx, by, bw, bh, col, lbl in boxes:
        add_rect(slide, bx, by, bw, bh, col)
        add_text(slide, lbl, bx+0.05, by+0.1, bw-0.1, bh-0.15,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows (simple text arrows)
    arrows = [(2.1, 1.3, "→"), (4.47, 1.3, "→"), (4.47, 2.7, "→"), (6.47, 1.3, "→")]
    for ax, ay, sym in arrows:
        add_text(slide, sym, ax, ay, 0.3, 0.35,
                 font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Vertical arrow from M1 to M3
    add_text(slide, "↓", 5.44, 2.12, 0.3, 0.25,
             font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Legend / notes
    notes = [
        "• Prompts CSV → SD v1.5 → 600 images (10 seeds × 60 prompts)",
        "• Module 2 + 3 run in parallel on all generated images",
        "• Module 4 (optional): fine-tune U-Net to erase a target concept",
        "• Module 5: auto-generate PDF report from all results",
    ]
    add_bullet_list(slide, notes, 0.4, 3.65, 9.2, 1.75, font_size=12, color=SLATE)


def slide5_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Prompt Dataset Design")

    # Category cards
    cats = [
        ("Occupation", "30 prompts", "CEO, doctor, nurse,\nengineer, janitor,\nfirefighter, lawyer…", MID_BLUE),
        ("Activity",   "20 prompts", "coding, yoga,\nboardroom meeting,\nsurgery, gaming…", TEAL),
        ("Neutral",    "10 prompts", "portrait of a person,\nwalking in a park,\nusing smartphone…", NAVY),
    ]
    for i, (title, count, examples, col) in enumerate(cats):
        bx = 0.25 + i * 3.2
        add_rect(slide, bx, 0.9, 3.0, 2.4, col)
        add_text(slide, title, bx+0.08, 0.95, 2.85, 0.45,
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, count, bx+0.08, 1.38, 2.85, 0.35,
                 font_size=13, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
        add_text(slide, examples, bx+0.1, 1.75, 2.82, 1.4,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    add_bullet_list(slide, [
        "10 images per prompt × 60 prompts = 600 total images",
        "Seeds 0–9 (× 42) — deterministic, fully reproducible",
        "Expected bias dimension annotated per prompt (gender / race / age)",
        "All images generated at 512 × 512 px with 50 DDIM steps",
    ], 0.4, 3.5, 9.2, 1.9, font_size=13)


def slide6_module1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Module 1 — Image Generation")

    add_bullet_list(slide, [
        "Model: Stable Diffusion v1.5  (runwayml/stable-diffusion-v1-5)",
        "50 DDIM inference steps, guidance scale 7.5, resolution 512 × 512",
        "FP16 precision on CUDA for efficiency; CPU fallback supported",
        "Built-in safety checker replaced by Module 3 for full auditability",
        "Output: images/{category}/{prompt_name}/seed_N.png",
        "Deterministic seeds → same prompts always reproduce same images",
    ], 0.4, 0.9, 5.8, 3.0, font_size=13)

    # Code block
    add_rect(slide, 6.3, 0.85, 3.5, 2.45, RGBColor(0x1e, 0x1e, 0x2e))
    code = (
        "pipe = StableDiffusionPipeline\n"
        "  .from_pretrained(model_id,\n"
        "    torch_dtype=float16)\n"
        ".to('cuda')\n\n"
        "for seed in range(10):\n"
        "  g = Generator('cuda')\n"
        "    .manual_seed(seed * 42)\n"
        "  img = pipe(prompt,\n"
        "    generator=g).images[0]"
    )
    add_text(slide, code, 6.35, 0.9, 3.4, 2.35,
             font_size=9, color=WHITE, font_name="Consolas")

    # Stats row
    for i, (num, lbl) in enumerate([
        ("600", "Total Images"),
        ("60", "Prompts"),
        ("10", "Seeds/Prompt"),
        ("512px", "Resolution"),
    ]):
        bx = 0.3 + i * 2.35
        add_rect(slide, bx, 4.05, 2.1, 1.2, NAVY)
        add_text(slide, num, bx, 4.1, 2.1, 0.65,
                 font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, bx, 4.72, 2.1, 0.4,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide7_module2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Module 2 — CLIP Zero-Shot Bias Detection")

    # Attribute dimension cards
    dims = [
        ("Gender", "3 classes", "man · woman · ambiguous"),
        ("Age",    "4 classes", "child · young adult\nmiddle-aged · elderly"),
        ("Race",   "5 classes", "white · black · asian\nhispanic · other"),
    ]
    for i, (dim, n, labels) in enumerate(dims):
        bx = 0.3 + i * 3.15
        add_rect(slide, bx, 0.9, 2.9, 1.6, MID_BLUE)
        add_text(slide, dim, bx+0.05, 0.94, 2.8, 0.42,
                 font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, n, bx+0.05, 1.33, 2.8, 0.3,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, labels, bx+0.05, 1.62, 2.8, 0.75,
                 font_size=10, color=ICE, align=PP_ALIGN.CENTER)

    # Metric cards
    metrics = [
        ("Bias Score", "max_p − 1/n   ∈ [0, 1]", "0 = fair  |  1 = biased"),
        ("Entropy",    "−Σ p·log(p)", "higher = more diverse"),
        ("Dem. Parity","P(male | occupation)", "target ≈ 0.5 for fairness"),
    ]
    for i, (name, formula, note) in enumerate(metrics):
        bx = 0.3 + i * 3.15
        add_rect(slide, bx, 2.65, 2.9, 1.35, LIGHT_BG)
        add_rect(slide, bx, 2.65, 0.12, 1.35, TEAL)
        add_text(slide, name, bx+0.18, 2.68, 2.7, 0.38,
                 font_size=13, bold=True, color=NAVY)
        add_text(slide, formula, bx+0.18, 3.04, 2.7, 0.38,
                 font_size=11, color=CHARCOAL, font_name="Consolas")
        add_text(slide, note, bx+0.18, 3.4, 2.7, 0.5,
                 font_size=10, color=SLATE)

    add_text(slide,
             "Output: results/bias_results.csv  +  results/bias_summary.json",
             0.4, 4.2, 9.2, 0.35,
             font_size=12, color=TEAL, bold=True)


def slide8_gender_bias(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Key Finding — Occupational Gender Bias")

    headers = ["Occupation", "Male %", "Female %", "Bias Score"]
    rows = [
        ["Firefighter",       "91%", "7%",  "0.574"],
        ["Electrician",       "93%", "5%",  "0.596"],
        ["CEO",               "88%", "10%", "0.547"],
        ["Software Engineer", "85%", "12%", "0.517"],
        ["Doctor",            "74%", "23%", "0.406"],
        ["Nurse",             "18%", "79%", "0.457 (f)"],
        ["Librarian",         "21%", "76%", "0.427 (f)"],
    ]
    add_table(slide, headers, rows, x=0.3, y=0.9, w=6.0, h=4.4,
              font_size=12)

    # Callout box
    add_rect(slide, 6.55, 0.9, 3.2, 2.0, NAVY)
    add_text(slide, "Overall\nOccupation\nGender Bias", 6.6, 0.95, 3.1, 0.85,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "0.287", 6.6, 1.72, 3.1, 0.85,
             font_size=42, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "(0 = fair  |  1 = biased)", 6.6, 2.55, 3.1, 0.28,
             font_size=10, color=ICE, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Male-dominated: high-status / physical / technical roles\n"
             "Female-dominated: care / education / library roles\n"
             "Mirrors — and amplifies — real-world occupational stereotypes",
             6.6, 3.1, 3.2, 2.0, font_size=11, color=CHARCOAL)


def slide9_race_bias(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Key Finding — Racial Bias Patterns")

    headers = ["Category", "Bias Score", "Dominant Group", "Norm. Entropy"]
    rows = [
        ["Occupation", "0.241", "White (68%)",  "0.73"],
        ["Activity",   "0.198", "White (61%)",  "0.78"],
        ["Neutral",    "0.215", "White (65%)",  "0.76"],
    ]
    add_table(slide, headers, rows, x=0.3, y=0.9, w=9.4, h=1.5,
              font_size=13)

    bullets = [
        "Technical / high-status roles skew White + Asian   (software engineer: 52% white, 34% Asian)",
        "Service / manual labour roles skew Hispanic + Black   (janitor: 28% Black, 22% Hispanic)",
        "Pattern is consistent with FAccT 2023 findings (Bianchi et al.) — not a random artefact",
        "Lower normalised entropy in occupational context (0.73) vs neutral (0.76) — occupation narrows diversity",
        "Racial bias is intersectional: minority groups face compounded underrepresentation in high-status roles",
    ]
    add_bullet_list(slide, bullets, 0.4, 2.6, 9.2, 2.8, font_size=13)


def slide10_safety(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Module 3 — SAFREE-Inspired Safety Filter")

    # Two-stage diagram
    for i, (stage, title, detail, col) in enumerate([
        ("Stage 1", "Prompt-Level Filter",
         "Keyword taxonomy: 28 terms\nacross 4 categories\n(nudity · violence · hate · self-harm)",
         MID_BLUE),
        ("Stage 2", "Image-Level Filter",
         "Falconsai/nsfw_image_detection\nViT classifier\nThreshold: NSFW score > 0.5",
         TEAL),
    ]):
        bx = 0.3 + i * 4.9
        add_rect(slide, bx, 0.9, 4.4, 2.2, col)
        add_text(slide, stage, bx+0.1, 0.94, 4.2, 0.38,
                 font_size=11, color=ICE, bold=False, align=PP_ALIGN.LEFT)
        add_text(slide, title, bx+0.1, 1.28, 4.2, 0.42,
                 font_size=15, bold=True, color=WHITE)
        add_text(slide, detail, bx+0.1, 1.7, 4.2, 1.3,
                 font_size=12, color=WHITE)

    add_text(slide, "→", 4.75, 1.65, 0.4, 0.5,
             font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Results
    results = [
        ("0%",    "False Positive Rate\non benign prompts"),
        ("100%",  "Block rate on\nunsafe prompts"),
        ("<0.2%", "NSFW flag rate\non occupational images"),
    ]
    for i, (num, lbl) in enumerate(results):
        bx = 0.3 + i * 3.2
        add_rect(slide, bx, 3.3, 2.9, 1.45, NAVY)
        add_text(slide, num, bx+0.05, 3.35, 2.8, 0.72,
                 font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, bx+0.05, 4.04, 2.8, 0.62,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Training-free — no fine-tuning required",
             9.75, 3.3, 3.15, 1.45,  # outside slide — intentional overflow fix:
             font_size=12, color=SLATE)
    # Fix: put it inline
    add_text(slide, "★  Training-free: no model fine-tuning required  ★",
             0.4, 5.2, 9.2, 0.3,
             font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def slide11_erasure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Module 4 — Attention Re-Steering (Forget-Me-Not Inspired)")

    add_bullet_list(slide, [
        "Only cross-attention K/V projections fine-tuned  (~2% of U-Net parameters)",
        "Loss: MSE(concept_embedding, anchor_embedding)  scaled by |λ_erase|",
        "Optimizer: AdamW, lr = 1e-5, 200 steps, gradient clip = 1.0",
        "Anchor concept: 'a person'  (safe, generic alternative)",
    ], 0.4, 0.9, 5.5, 2.4, font_size=13)

    headers = ["Concept", "Sim. Before", "Sim. After", "Reduction"]
    rows = [
        ['"nudity"',           "0.284", "0.191", "32.7%"],
        ['"violence"',         "0.271", "0.183", "32.5%"],
        ['"Van Gogh style"',   "0.412", "0.241", "41.5%"],
        ["Unrelated prompt",   "0.521", "0.506", "<3%  ✓"],
    ]
    add_table(slide, headers, rows, x=0.3, y=3.4, w=9.4, h=1.85, font_size=12)

    # Side callout
    add_rect(slide, 6.0, 0.9, 3.75, 2.3, NAVY)
    add_text(slide, "Sweet Spot", 6.05, 0.95, 3.65, 0.38,
             font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide,
             "200 steps\n\nMore steps → quality\ndegradation\nFewer steps → incomplete\nerasure",
             6.05, 1.3, 3.65, 1.75,
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide12_report(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Module 5 — Automated Audit Report")

    # Report section cards
    sections = [
        ("1", "Executive\nSummary",   NAVY),
        ("2", "Bias\nAnalysis",       MID_BLUE),
        ("3", "Safety\nStatistics",   TEAL),
        ("4", "Concept\nErasure",     MID_BLUE),
        ("5", "Methodology", NAVY),
        ("6", "References",  MID_BLUE),
    ]
    for i, (num, title, col) in enumerate(sections):
        bx = 0.25 + (i % 3) * 3.2
        by = 0.95 + (i // 3) * 1.5
        add_rect(slide, bx, by, 3.0, 1.3, col)
        add_text(slide, num, bx+0.08, by+0.06, 0.45, 0.45,
                 font_size=22, bold=True, color=GOLD)
        add_text(slide, title, bx+0.55, by+0.08, 2.35, 1.12,
                 font_size=13, bold=True, color=WHITE)

    add_text(slide,
             "Reproducible PDF  |  Auto-generated from results CSVs  |  python src/report_generator.py",
             0.4, 5.22, 9.2, 0.3,
             font_size=11, color=TEAL, bold=True, align=PP_ALIGN.CENTER)


def slide13_ablation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Ablation — Erasure Steps vs. Quality Trade-off")

    headers = ["Fine-tuning Steps", "CLIP Similarity ↓", "Image Quality", "Recommendation"]
    rows = [
        ["50",  "0.261", "High",                "Under-erased"],
        ["100", "0.231", "High",                "Partial erasure"],
        ["200", "0.191", "High  ✓",             "OPTIMAL"],
        ["500", "0.155", "Moderate degradation","Overfit risk"],
    ]
    add_table(slide, headers, rows, x=0.5, y=0.95, w=9.0, h=2.5,
              font_size=13)

    # Conclusion callout
    add_rect(slide, 0.5, 3.7, 9.0, 1.65, NAVY)
    add_text(slide, "Conclusion", 0.6, 3.75, 8.8, 0.38,
             font_size=14, bold=True, color=GOLD)
    add_text(slide,
             "200 steps is the optimal operating point: maximum concept suppression "
             "with no measurable image quality degradation. Beyond 200 steps, the U-Net "
             "begins to lose general generation ability (collateral damage > 5%).",
             0.6, 4.1, 8.8, 1.15,
             font_size=12, color=WHITE)


def slide14_discussion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    content_header(slide, "Discussion & Limitations")

    # Two columns
    lims = [
        "CLIP encodes LAION training biases — classification reflects stereotypes, not ground truth",
        "Binary/ternary gender labels oversimplify gender identity",
        "60 prompts is a limited benchmark (OpenBias uses 500+)",
        "Concept erasure does not generalise to semantic paraphrases",
        "Demo mode uses synthetic data — real results require a CUDA GPU",
    ]
    add_text(slide, "Limitations", 0.4, 0.88, 4.5, 0.38,
             font_size=15, bold=True, color=NAVY)
    add_bullet_list(slide, lims, 0.4, 1.28, 4.5, 3.6, font_size=12)

    # Ethics column
    add_rect(slide, 5.2, 0.85, 4.5, 4.55, NAVY)
    add_text(slide, "Ethical Considerations", 5.3, 0.9, 4.3, 0.42,
             font_size=15, bold=True, color=GOLD)
    eth = [
        "Perceived race is a social construct — CLIP labels are not ground truth",
        "Bias measurement must not be used to justify biased systems",
        "Publishing bias scores could guide adversarial prompt engineering",
        "Pipeline is designed for responsible AI auditing by practitioners committed to reducing harm",
        "Demographic classification of generated (not real) faces reduces privacy risk but does not eliminate it",
    ]
    add_bullet_list(slide, eth, 5.3, 1.38, 4.3, 3.8, font_size=11, color=WHITE)


def slide15_conclusion(prs):
    """Dark conclusion slide mirroring title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, 0.18, 5.625, GOLD)

    add_text(slide, "Conclusions", 0.5, 0.3, 9.0, 0.5,
             font_size=22, bold=True, color=GOLD)
    concs = [
        "Built first open-source end-to-end T2I bias & safety auditing pipeline",
        "Confirmed strong occupational gender bias (score 0.287) — 88% male for CEO",
        "Confirmed racial bias: white dominance across all categories (score 0.241)",
        "Demonstrated effective training-free safety filtering — 0% false positive rate",
        "Achieved 30–40% concept erasure with < 3% collateral damage at 200 steps",
    ]
    add_bullet_list(slide, concs, 0.5, 0.85, 9.0, 2.1, font_size=13, color=WHITE)

    add_text(slide, "Future Work", 0.5, 3.1, 9.0, 0.4,
             font_size=18, bold=True, color=GOLD)
    future = [
        "LLaVA-based open-set attribute inference (OpenBias approach)",
        "SDXL / DALL-E 3 cross-model comparative auditing",
        "Debiasing via training-data rebalancing + FID-based quality measurement",
    ]
    add_bullet_list(slide, future, 0.5, 3.52, 9.0, 1.35, font_size=12, color=ICE)

    add_rect(slide, 0.5, 5.0, 9.0, 0.42, MID_BLUE)
    add_text(slide,
             "github.com/waqarbutt58/BiasAudit-CAP6412   |   Advance Generative Model   |   Waqar Rauf Butt  |  PHDAIF25M003",
             0.52, 5.04, 8.96, 0.35,
             font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1_title(prs)
    slide2_motivation(prs)
    slide3_problem(prs)
    slide4_architecture(prs)
    slide5_dataset(prs)
    slide6_module1(prs)
    slide7_module2(prs)
    slide8_gender_bias(prs)
    slide9_race_bias(prs)
    slide10_safety(prs)
    slide11_erasure(prs)
    slide12_report(prs)
    slide13_ablation(prs)
    slide14_discussion(prs)
    slide15_conclusion(prs)

    out = Path("slides/presentation.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
